# converter_api.py
#
# QuadraConverter - production-oriented conversion API
#
# Supports:
#   - Office -> PDF
#   - HTML -> PDF
#   - PDF -> Word
#   - PDF -> Excel
#   - PDF -> PowerPoint
#   - PDF -> PDF/A
#   - PDF Unlock
#   - PDF Protect
#   - PDF Translate
#   - Send converted file by email
#
# PDF -> Word:
#   Native PDF text extraction
#   Coordinate-aware line reconstruction
#   Font-size estimation
#   Scanned-PDF OCR fallback
#
# PDF -> Excel:
#   Native table detection
#   Line/coordinate extraction
#   Column clustering
#   Scanned-PDF OCR fallback
#
# IMPORTANT:
# This file intentionally contains ONLY ONE /convert route.
#
# Required Python packages:
#   fastapi
#   uvicorn
#   python-multipart
#   PyMuPDF
#   pdfplumber
#   python-docx
#   openpyxl
#   python-pptx
#   Pillow
#   pytesseract
#   resend
#
# Required system programs:
#   LibreOffice
#   qpdf
#   Ghostscript
#   Tesseract OCR
#
# ============================================================

import hashlib
import math
import httpx
import os
import re
import json
import shutil
import subprocess
import tempfile
import urllib.request
from pathlib import Path
from io import BytesIO
from typing import Any
from urllib.parse import quote

from fastapi import (
    FastAPI,
    File,
    Form,
    HTTPException,
    UploadFile,
)

from fastapi.middleware.cors import CORSMiddleware

from fastapi.responses import (
    FileResponse,
    JSONResponse,
)

from starlette.background import BackgroundTask

from dotenv import load_dotenv

load_dotenv(
    Path(__file__).resolve().parent / ".env"
)
# ============================================================
# CONFIGURATION
# ============================================================

APP_NAME = "QuadraConverter Conversion API"
APP_VERSION = "3.0.0"


MAX_FILE_BYTES = int(
    os.getenv(
        "MAX_FILE_BYTES",
        str(100 * 1024 * 1024),
    )
)


OCR_DPI = int(
    os.getenv(
        "OCR_DPI",
        "250",
    )
)


OCR_MIN_TEXT_CHARS = int(
    os.getenv(
        "OCR_MIN_TEXT_CHARS",
        "40",
    )
)


CONVERSION_TIMEOUT = int(
    os.getenv(
        "CONVERSION_TIMEOUT",
        "600",
    )
)


MAX_PDF_PAGES = int(
    os.getenv(
        "MAX_PDF_PAGES",
        "500",
    )
)

SIGNING_KEY_FILE = os.getenv("QUADRA_SIGNING_KEY_FILE", "")
SIGNING_CERT_FILE = os.getenv("QUADRA_SIGNING_CERT_FILE", "")
SIGNING_CHAIN_FILE = os.getenv("QUADRA_SIGNING_CHAIN_FILE", "")
SIGNING_KEY_PASSPHRASE = os.getenv("QUADRA_SIGNING_KEY_PASSPHRASE", "")
SIGNING_TSA_URL = os.getenv("QUADRA_TSA_URL", "")
SIGNING_SUBJECT = os.getenv("QUADRA_SIGNING_SUBJECT", "QuadraConverter Production Signer")


CORS_ORIGINS = [
    origin.strip()
    for origin in os.getenv(
        "CORS_ORIGINS",
        "*",
    ).split(",")
    if origin.strip()
]


ALLOW_CREDENTIALS = (
    "*" not in CORS_ORIGINS
)


# ============================================================
# SUPPORTED FILE TYPES
# ============================================================

ALLOWED_PDF = {
    ".pdf",
}


ALLOWED_OFFICE = {
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".odt",
    ".ods",
    ".odp",
    ".rtf",
}


ALLOWED_HTML = {
    ".html",
    ".htm",
    ".xhtml",
}


# ============================================================
# FASTAPI APP
# ============================================================

app = FastAPI(
    title=APP_NAME,
    version=APP_VERSION,
)


app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_credentials=ALLOW_CREDENTIALS,
    allow_methods=[
        "GET",
        "POST",
        "OPTIONS",
    ],
    allow_headers=["*"],
    expose_headers=[
        "Content-Disposition",
        "X-Converted-Filename",
        "X-Conversion-Engine",
    ],
)


@app.get("/health")
async def health():
    return {
        "ok": True,
        "service": "quadraconverter-signing-api",
        "version": APP_VERSION,
        "pki_configured": bool(SIGNING_KEY_FILE and SIGNING_CERT_FILE),
        "tsa_configured": bool(SIGNING_TSA_URL),
    }


# ============================================================
# GENERAL HELPERS
# ============================================================


def cleanup(path: Path | None):
    if path is None:
        return

    try:
        shutil.rmtree(
            path,
            ignore_errors=True,
        )
    except Exception:
        pass


def binary_path(
    *names: str,
) -> str:

    for name in names:

        found = shutil.which(name)

        if found:
            return found

    raise RuntimeError(
        "Required converter is not installed: "
        + ", ".join(names)
    )


def command_exists(
    *names: str,
) -> bool:

    for name in names:

        if shutil.which(name):
            return True

    return False


def safe_filename(
    filename: str | None,
    fallback: str = "converted-file",
) -> str:

    name = (
        Path(filename or fallback).name
    )

    name = re.sub(
        r"[^A-Za-z0-9._() \-]+",
        "_",
        name,
    )

    name = name.strip(
        " ."
    )

    return name or fallback


def validate_file_size(
    path: Path,
):

    if not path.exists():
        raise HTTPException(
            status_code=422,
            detail="Uploaded file was not saved.",
        )

    size = path.stat().st_size

    if size <= 0:
        raise HTTPException(
            status_code=400,
            detail="The uploaded file is empty.",
        )

    if size > MAX_FILE_BYTES:
        raise HTTPException(
            status_code=413,
            detail=(
                f"File exceeds the "
                f"{MAX_FILE_BYTES // (1024 * 1024)} MB "
                f"conversion limit."
            ),
        )


def save_upload(
    upload: UploadFile,
    work: Path,
    allowed: set[str],
) -> Path:

    original_name = (
        upload.filename
        or "input"
    )

    suffix = (
        Path(original_name)
        .suffix
        .lower()
    )

    if suffix not in allowed:

        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported file type: "
                f"{suffix or 'unknown'}"
            ),
        )

    filename = safe_filename(
        original_name,
        f"input{suffix}",
    )

    source = (
        work / filename
    )

    size = 0

    try:

        with source.open(
            "wb"
        ) as output:

            while True:

                chunk = upload.file.read(
                    1024 * 1024
                )

                if not chunk:
                    break

                size += len(chunk)

                if size > MAX_FILE_BYTES:

                    raise HTTPException(
                        status_code=413,
                        detail=(
                            f"File exceeds the "
                            f"{MAX_FILE_BYTES // (1024 * 1024)} MB "
                            f"conversion limit."
                        ),
                    )

                output.write(
                    chunk
                )

    except HTTPException:
        raise

    except Exception as exc:

        raise HTTPException(
            status_code=400,
            detail=(
                f"Could not save uploaded file: "
                f"{exc}"
            ),
        )

    validate_file_size(
        source
    )

    return source


def run_checked(
    args: list[str],
    timeout: int = CONVERSION_TIMEOUT,
):

    try:

        process = subprocess.run(
            args,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=timeout,
            check=False,
        )

    except FileNotFoundError as exc:

        raise HTTPException(
            status_code=503,
            detail=(
                "Required conversion engine "
                f"is not installed: {exc}"
            ),
        )

    except subprocess.TimeoutExpired:

        raise HTTPException(
            status_code=504,
            detail=(
                "Conversion timed out. "
                "Try a smaller or simpler file."
            ),
        )

    stdout = (
        process.stdout or ""
    ).strip()

    stderr = (
        process.stderr or ""
    ).strip()

    if process.returncode != 0:

        detail = (
            stderr
            or stdout
            or "Conversion engine failed."
        )

        raise HTTPException(
            status_code=422,
            detail=detail[-5000:],
        )

    return process


def verify_output(
    output: Path,
    description: str,
):

    if not output.exists():

        raise HTTPException(
            status_code=422,
            detail=(
                f"{description} did not "
                "produce an output file."
            ),
        )

    if output.stat().st_size <= 0:

        raise HTTPException(
            status_code=422,
            detail=(
                f"{description} produced "
                "an empty output file."
            ),
        )

    return output


def file_response(
    output: Path,
    media_type: str,
    work: Path,
    engine: str,
):
    verify_output(
        output,
        "Conversion",
    )

    filename = output.name

    encoded_filename = quote(
        filename,
        safe="",
    )

    return FileResponse(
        path=str(output),
        media_type=media_type,
        filename=filename,
        headers={
            "Content-Disposition": (
                f'attachment; filename="{filename}"; '
                f"filename*=UTF-8''{encoded_filename}"
            ),
            "X-Converted-Filename": filename,
            "X-Conversion-Engine": engine,
            "Cache-Control": "no-store",
        },
        background=BackgroundTask(
            cleanup,
            work,
        ),
    )


# ============================================================
# OCR
# ============================================================


def normalize_ocr_language(
    language: str,
) -> str:

    value = (
        language or "eng"
    ).strip().lower()

    aliases = {
        "english": "eng",
        "en": "eng",

        "hindi": "hin",
        "hi": "hin",

        "tamil": "tam",
        "ta": "tam",

        "telugu": "tel",
        "te": "tel",

        "malayalam": "mal",
        "ml": "mal",

        "kannada": "kan",
        "kn": "kan",

        "marathi": "mar",
        "mr": "mar",

        "bengali": "ben",
        "bn": "ben",

        "gujarati": "guj",
        "gu": "guj",

        "punjabi": "pan",
        "pa": "pan",

        "odia": "ori",
        "or": "ori",

        "urdu": "urd",
        "ur": "urd",

        "sanskrit": "san",
        "sa": "san",
    }

    parts = []

    for item in re.split(
        r"[+,;\s]+",
        value,
    ):
        item = item.strip()

        if not item:
            continue

        parts.append(
            aliases.get(
                item,
                item,
            )
        )

    return "+".join(
        dict.fromkeys(parts)
    ) or "eng"


def tesseract_available() -> bool:

    return (
        shutil.which(
            "tesseract"
        )
        is not None
    )


def installed_tesseract_languages() -> set[str]:

    if not tesseract_available():
        return set()

    try:

        process = subprocess.run(
            [
                "tesseract",
                "--list-langs",
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            timeout=20,
            check=False,
        )

        lines = (
            process.stdout or ""
        ).splitlines()

        result = set()

        for line in lines:

            value = line.strip()

            if (
                value
                and not value.startswith(
                    "List of available languages"
                )
            ):
                result.add(
                    value
                )

        return result

    except Exception:

        return set()


def render_page_image(
    page,
    dpi: int = OCR_DPI,
):

    import fitz

    scale = dpi / 72.0

    matrix = fitz.Matrix(
        scale,
        scale,
    )

    pixmap = page.get_pixmap(
        matrix=matrix,
        alpha=False,
    )

    return pixmap, scale


def run_ocr(
    page,
    language: str = "eng",
):

    if not tesseract_available():

        raise HTTPException(
            status_code=503,
            detail=(
                "This PDF appears to be scanned, "
                "but Tesseract OCR is not installed "
                "on the conversion server."
            ),
        )

    try:

        import pytesseract
        from PIL import Image

    except ImportError as exc:

        raise HTTPException(
            status_code=503,
            detail=(
                "OCR dependencies are missing: "
                f"{exc}"
            ),
        )

    available = (
        installed_tesseract_languages()
    )

    requested_languages = (
        language.split("+")
    )

    usable_languages = [
        item
        for item in requested_languages
        if not available
        or item in available
    ]

    if not usable_languages:

        raise HTTPException(
            status_code=422,
            detail=(
                "The requested OCR language "
                f"'{language}' is not installed "
                "on the conversion server."
            ),
        )

    ocr_language = "+".join(
        usable_languages
    )

    pixmap, scale = (
        render_page_image(
            page
        )
    )

    image = Image.open(
        BytesIO(
            pixmap.tobytes(
                "png"
            )
        )
    )

    try:

        data = (
            pytesseract.image_to_data(
                image,
                lang=ocr_language,
                config="--oem 3 --psm 3",
                output_type=(
                    pytesseract.Output.DICT
                ),
            )
        )

    except Exception as exc:

        raise HTTPException(
            status_code=422,
            detail=(
                f"OCR failed: {exc}"
            ),
        )

    words = []

    count = len(
        data.get(
            "text",
            [],
        )
    )

    for index in range(count):

        text = (
            data["text"][index]
            or ""
        ).strip()

        if not text:
            continue

        try:

            confidence = float(
                data["conf"][index]
            )

        except (
            ValueError,
            TypeError,
        ):

            confidence = -1

        if confidence < 10:
            continue

        x = int(
            data["left"][index]
        )

        y = int(
            data["top"][index]
        )

        width = int(
            data["width"][index]
        )

        height = int(
            data["height"][index]
        )

        words.append(
            {
                "text": text,
                "x0": x / scale,
                "x1": (
                    x + width
                ) / scale,
                "top": y / scale,
                "bottom": (
                    y + height
                ) / scale,
                "height": max(
                    1.0,
                    height / scale,
                ),
                "confidence": confidence,
                "source": "ocr",
            }
        )

    return words


# ============================================================
# PDF NATIVE TEXT EXTRACTION
# ============================================================


def extract_pdf_words(
    page,
    language: str = "eng",
):

    words = []

    try:

        raw_words = (
            page.get_text(
                "words",
                sort=True,
            )
            or []
        )

    except Exception:

        raw_words = []

    for item in raw_words:

        if len(item) < 5:
            continue

        text = str(
            item[4]
        ).strip()

        if not text:
            continue

        x0 = float(
            item[0]
        )

        top = float(
            item[1]
        )

        x1 = float(
            item[2]
        )

        bottom = float(
            item[3]
        )

        words.append(
            {
                "text": text,
                "x0": x0,
                "x1": x1,
                "top": top,
                "bottom": bottom,
                "height": max(
                    1.0,
                    bottom - top,
                ),
                "source": "pdf",
            }
        )

    text_length = sum(
        len(
            word["text"]
        )
        for word in words
    )

    if (
        text_length
        < OCR_MIN_TEXT_CHARS
    ):

        try:

            ocr_words = run_ocr(
                page,
                language,
            )

            if ocr_words:
                return ocr_words

        except HTTPException:
            raise

        except Exception:
            pass

    return words


# ============================================================
# LINE RECONSTRUCTION
# ============================================================


def group_words_into_lines(
    words: list[dict[str, Any]],
    tolerance: float = 4.0,
):

    if not words:
        return []

    ordered = sorted(
        words,
        key=lambda word: (
            word["top"],
            word["x0"],
        ),
    )

    lines = []

    for word in ordered:

        center_y = (
            word["top"]
            + word["bottom"]
        ) / 2.0

        selected = None

        for line in reversed(
            lines[-8:]
        ):

            line_height = max(
                1.0,
                line["height"],
            )

            tolerance_value = max(
                tolerance,
                line_height * 0.55,
                word["height"] * 0.55,
            )

            if abs(
                center_y
                - line["center_y"]
            ) <= tolerance_value:

                selected = line
                break

        if selected is None:

            selected = {
                "center_y": center_y,
                "height": word[
                    "height"
                ],
                "words": [],
            }

            lines.append(
                selected
            )

        selected[
            "words"
        ].append(
            word
        )

        selected["height"] = max(
            selected["height"],
            word["height"],
        )

        selected["center_y"] = (
            selected["center_y"]
            + center_y
        ) / 2.0

    result = []

    for line in lines:

        line["words"].sort(
            key=lambda word:
            word["x0"]
        )

        if not line["words"]:
            continue

        line["text"] = " ".join(
            word["text"]
            for word in line[
                "words"
            ]
        ).strip()

        line["x0"] = min(
            word["x0"]
            for word in line[
                "words"
            ]
        )

        line["x1"] = max(
            word["x1"]
            for word in line[
                "words"
            ]
        )

        line["top"] = min(
            word["top"]
            for word in line[
                "words"
            ]
        )

        line["bottom"] = max(
            word["bottom"]
            for word in line[
                "words"
            ]
        )

        result.append(
            line
        )

    result.sort(
        key=lambda line: (
            line["top"],
            line["x0"],
        )
    )

    return result


def line_font_size(
    line: dict[str, Any],
) -> float:

    heights = [
        word["height"]
        for word in line[
            "words"
        ]
        if word.get("height")
    ]

    if not heights:
        return 10.0

    average = (
        sum(heights)
        / len(heights)
    )

    return max(
        6.0,
        min(
            32.0,
            average * 0.88,
        ),
    )


# ============================================================
# PDF → WORD
# ============================================================


def _pdf_font_name(name: str | None) -> str:
    """Map common PDF font names to fonts that are usually available in Word."""
    value = (name or "Arial").split("+")[-1].strip()
    lowered = value.lower()
    if "times" in lowered or "roman" in lowered:
        return "Times New Roman"
    if "courier" in lowered or "mono" in lowered:
        return "Courier New"
    if "arial" in lowered or "helvetica" in lowered or "calibri" in lowered:
        return "Arial"
    if "dejavu" in lowered:
        return "DejaVu Sans"
    return "Arial"


def _pdf_color_rgb(value: int | None):
    from docx.shared import RGBColor

    packed = int(value or 0)
    return RGBColor(
        (packed >> 16) & 0xFF,
        (packed >> 8) & 0xFF,
        packed & 0xFF,
    )


def _configure_docx_section(section, page):
    from docx.shared import Inches

    width_in = max(1.0, float(page.rect.width) / 72.0)
    height_in = max(1.0, float(page.rect.height) / 72.0)

    section.page_width = Inches(width_in)
    section.page_height = Inches(height_in)
    section.top_margin = Inches(0.25)
    section.bottom_margin = Inches(0.25)
    section.left_margin = Inches(0.25)
    section.right_margin = Inches(0.25)
    section.header_distance = Inches(0.1)
    section.footer_distance = Inches(0.1)


def _add_docx_text_block(document, block, page_width_pt: float):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt
    from docx.oxml.ns import qn

    bbox = block.get("bbox") or (0, 0, page_width_pt, 0)
    x0, y0, x1, y1 = [float(v) for v in bbox]
    paragraph = document.add_paragraph()
    pf = paragraph.paragraph_format

    # Approximate the original PDF horizontal position.
    pf.left_indent = Inches(max(0.0, min((x0 / 72.0) - 0.25, max(0.0, page_width_pt / 72.0 - 0.5))))
    pf.space_before = Pt(0)
    pf.space_after = Pt(0)
    pf.line_spacing = 1.0

    lines = block.get("lines") or []
    for line_index, line in enumerate(lines):
        if line_index:
            paragraph.add_run().add_break()

        spans = line.get("spans") or []
        for span in spans:
            value = str(span.get("text") or "")
            if not value:
                continue

            run = paragraph.add_run(value)
            run.font.size = Pt(max(6.0, min(72.0, float(span.get("size") or 10.0))))
            font_name = _pdf_font_name(span.get("font"))
            run.font.name = font_name
            try:
                run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
            except Exception:
                pass

            flags = int(span.get("flags") or 0)
            run.bold = bool(flags & 16)
            run.italic = bool(flags & 2)

            try:
                run.font.color.rgb = _pdf_color_rgb(span.get("color"))
            except Exception:
                pass

    # A paragraph containing only whitespace is not useful in DOCX.
    if not paragraph.text.strip():
        document._body._body.remove(paragraph._p)
        return False

    return True


def _add_docx_image_block(document, block, page_width_pt: float):
    from docx.shared import Inches

    image_bytes = block.get("image")
    if not image_bytes:
        return False

    bbox = block.get("bbox") or (0, 0, page_width_pt, page_width_pt)
    x0, y0, x1, y1 = [float(v) for v in bbox]
    width_pt = max(1.0, x1 - x0)
    max_width_in = max(1.0, page_width_pt / 72.0 - 0.5)
    width_in = min(width_pt / 72.0, max_width_in)

    paragraph = document.add_paragraph()
    paragraph.paragraph_format.left_indent = Inches(max(0.0, (x0 / 72.0) - 0.25))
    paragraph.paragraph_format.space_before = Inches(max(0.0, y0 / 72.0) * 0.0)
    paragraph.paragraph_format.space_after = Inches(0)

    try:
        run = paragraph.add_run()
        run.add_picture(BytesIO(image_bytes), width=Inches(width_in))
        return True
    except Exception:
        document._body._body.remove(paragraph._p)
        return False


def pdf_to_docx(
    source: Path,
    output: Path,
    language: str = "eng",
):
    """Convert PDF pages into a real DOCX with editable native text and embedded images.

    Native PDF text is reconstructed from PyMuPDF blocks/spans so font size, font family,
    bold/italic state, color and approximate horizontal placement survive. Scanned pages
    fall back to the existing Tesseract word reconstruction.
    """
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.shared import Inches, Pt
    import fitz

    try:
        pdf = fitz.open(str(source))
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not open PDF: {exc}")

    try:
        if pdf.page_count == 0:
            raise HTTPException(status_code=422, detail="The PDF has no pages.")
        if pdf.page_count > MAX_PDF_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"PDF contains more than {MAX_PDF_PAGES} pages.",
            )

        document = Document()
        total_content = 0

        for page_index, page in enumerate(pdf):
            if page_index == 0:
                section = document.sections[0]
            else:
                section = document.add_section(WD_SECTION.NEW_PAGE)
            _configure_docx_section(section, page)

            page_dict = page.get_text("dict", sort=True) or {}
            blocks = page_dict.get("blocks") or []
            native_text = sum(
                len(str(span.get("text") or ""))
                for block in blocks
                if block.get("type") == 0
                for line in (block.get("lines") or [])
                for span in (line.get("spans") or [])
            )

            if native_text >= OCR_MIN_TEXT_CHARS:
                for block in blocks:
                    block_type = block.get("type")
                    if block_type == 0:
                        if _add_docx_text_block(document, block, float(page.rect.width)):
                            total_content += 1
                    elif block_type == 1:
                        if _add_docx_image_block(document, block, float(page.rect.width)):
                            total_content += 1
            else:
                # Scanned/image-only page: use OCR words so the DOCX remains editable.
                words = extract_pdf_words(page, language)
                lines = group_words_into_lines(words)
                previous_bottom = None

                for line in lines:
                    text = str(line.get("text") or "").strip()
                    if not text:
                        continue

                    paragraph = document.add_paragraph()
                    pf = paragraph.paragraph_format
                    left = max(0.0, (float(line.get("x0", 0)) / 72.0) - 0.25)
                    pf.left_indent = Inches(min(left, max(0.0, page.rect.width / 72.0 - 0.5)))
                    if previous_bottom is not None:
                        gap = max(0.0, float(line.get("top", 0)) - previous_bottom)
                        pf.space_before = Pt(min(30.0, gap))
                    pf.space_after = Pt(0)
                    pf.line_spacing = 1.0

                    font_size = line_font_size(line)
                    for index, word in enumerate(line.get("words") or []):
                        if index:
                            paragraph.add_run(" ")
                        run = paragraph.add_run(str(word.get("text") or ""))
                        run.font.size = Pt(font_size)
                    previous_bottom = float(line.get("bottom", 0))
                    total_content += 1

            # Preserve an entirely blank page.
            if not blocks and native_text == 0:
                try:
                    words = extract_pdf_words(page, language)
                except Exception:
                    words = []
                if not words:
                    document.add_paragraph("")

        if total_content == 0:
            raise HTTPException(
                status_code=422,
                detail="No readable content was found in this PDF. If it is scanned, make sure the required Tesseract language is installed.",
            )

        output.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(output))
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Could not create Word document: {exc}")
    finally:
        pdf.close()

    return verify_output(output, "PDF → Word")

def _pptx_font_name(name: str | None) -> str:
    value = (name or "Arial").split("+")[-1].strip()
    lowered = value.lower()
    if "times" in lowered or "roman" in lowered:
        return "Times New Roman"
    if "courier" in lowered or "mono" in lowered:
        return "Courier New"
    if "arial" in lowered or "helvetica" in lowered or "calibri" in lowered:
        return "Arial"
    if "dejavu" in lowered:
        return "DejaVu Sans"
    return "Arial"


def _pptx_color_rgb(value: int | None):
    from pptx.dml.color import RGBColor

    packed = int(value or 0)
    return RGBColor(
        (packed >> 16) & 0xFF,
        (packed >> 8) & 0xFF,
        packed & 0xFF,
    )


def _add_pptx_text_block(slide, block, page_rect, slide_width, slide_height):
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    bbox = block.get("bbox")
    if not bbox:
        return False

    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h

    x0, y0, x1, y1 = [float(v) for v in bbox]
    left = max(0, int(x0 * sx))
    top = max(0, int(y0 * sy))
    width = max(1, int((x1 - x0) * sx))
    height = max(1, int((y1 - y0) * sy))

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = block.get("lines") or []
    first_paragraph = True
    added = False

    for line in lines:
        paragraph = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0

        spans = line.get("spans") or []
        for span in spans:
            value = str(span.get("text") or "")
            if not value:
                continue
            run = paragraph.add_run()
            run.text = value
            run.font.name = _pptx_font_name(span.get("font"))
            run.font.size = Pt(max(6.0, min(96.0, float(span.get("size") or 10.0))))
            flags = int(span.get("flags") or 0)
            run.font.bold = bool(flags & 16)
            run.font.italic = bool(flags & 2)
            try:
                run.font.color.rgb = _pptx_color_rgb(span.get("color"))
            except Exception:
                pass
            added = True

    if not added:
        slide.shapes._spTree.remove(shape._element)
        return False

    # PDF text is generally left aligned; preserve obvious centered/right blocks.
    block_center = (x0 + x1) / 2.0
    page_center = page_w / 2.0
    if abs(block_center - page_center) <= page_w * 0.08:
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

    return True


def _add_pptx_image_block(slide, block, page_rect, slide_width, slide_height):
    from pptx.util import Inches

    image_bytes = block.get("image")
    bbox = block.get("bbox")
    if not image_bytes or not bbox:
        return False

    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h
    x0, y0, x1, y1 = [float(v) for v in bbox]

    left = int(x0 * sx)
    top = int(y0 * sy)
    width = max(1, int((x1 - x0) * sx))
    height = max(1, int((y1 - y0) * sy))

    try:
        slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)
        return True
    except Exception:
        return False


def _add_pptx_ocr_page(slide, page, language, slide_width, slide_height):
    """Add editable OCR text to a slide for scanned PDFs."""
    words = extract_pdf_words(page, language)
    lines = group_words_into_lines(words)
    if not lines:
        return False

    page_rect = page.rect
    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h
    added = False

    from pptx.util import Pt

    for line in lines:
        x0 = float(line.get("x0", 0))
        y0 = float(line.get("top", 0))
        x1 = float(line.get("x1", x0 + 10))
        y1 = float(line.get("bottom", y0 + 12))
        shape = slide.shapes.add_textbox(
            int(x0 * sx),
            int(y0 * sy),
            max(1, int((x1 - x0) * sx)),
            max(1, int((y1 - y0) * sy)),
        )
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        run = tf.paragraphs[0].add_run()
        run.text = str(line.get("text") or "")
        run.font.name = "Arial"
        run.font.size = Pt(line_font_size(line))
        added = True

    return added


def pdf_to_pptx(
    source: Path,
    output: Path,
    language: str = "eng",
):
    """Convert PDF to an editable PowerPoint.

    Native PDF text becomes PowerPoint text boxes with editable runs. Native PDF raster
    images become PowerPoint pictures. Scanned pages use OCR text boxes; a page screenshot
    is used only when a page contains no extractable content at all.
    """
    from pptx import Presentation
    from pptx.util import Inches
    import fitz

    pdf = _fitz_open(source)
    try:
        if pdf.page_count == 0:
            raise HTTPException(status_code=422, detail="The PDF has no pages.")
        if pdf.page_count > MAX_PDF_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"PDF contains more than {MAX_PDF_PAGES} pages.",
            )

        first_rect = pdf[0].rect
        page_ratio = float(first_rect.width) / max(1.0, float(first_rect.height))

        # Keep the longest slide dimension at 13.333 in so portrait PDFs do not create
        # invalid PowerPoint dimensions.
        if page_ratio >= 1:
            slide_width_in = 13.333
            slide_height_in = 13.333 / page_ratio
        else:
            slide_height_in = 13.333
            slide_width_in = 13.333 * page_ratio

        presentation = Presentation()
        presentation.slide_width = Inches(slide_width_in)
        presentation.slide_height = Inches(slide_height_in)
        blank_layout = presentation.slide_layouts[6]

        for page in pdf:
            slide = presentation.slides.add_slide(blank_layout)
            page_dict = page.get_text("dict", sort=True) or {}
            blocks = page_dict.get("blocks") or []

            native_text = sum(
                len(str(span.get("text") or ""))
                for block in blocks
                if block.get("type") == 0
                for line in (block.get("lines") or [])
                for span in (line.get("spans") or [])
            )

            added = False
            if native_text >= OCR_MIN_TEXT_CHARS:
                for block in blocks:
                    if block.get("type") == 0:
                        added = _add_pptx_text_block(
                            slide,
                            block,
                            page.rect,
                            presentation.slide_width,
                            presentation.slide_height,
                        ) or added
                    elif block.get("type") == 1:
                        added = _add_pptx_image_block(
                            slide,
                            block,
                            page.rect,
                            presentation.slide_width,
                            presentation.slide_height,
                        ) or added
            else:
                added = _add_pptx_ocr_page(
                    slide,
                    page,
                    language,
                    presentation.slide_width,
                    presentation.slide_height,
                )

            # If a PDF page truly contains no readable text/image blocks, retain it as
            # a visual page rather than returning an empty slide.
            if not added:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                slide.shapes.add_picture(
                    BytesIO(pixmap.tobytes("png")),
                    0,
                    0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )

        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output))
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Could not create PowerPoint: {exc}")
    finally:
        pdf.close()

    return verify_output(output, "PDF → PowerPoint")

def _fitz_open(
    source: Path,
):

    try:

        import fitz

        return fitz.open(
            str(source)
        )

    except Exception as exc:

        raise HTTPException(
            status_code=422,
            detail=(
                f"Could not read PDF: "
                f"{exc}"
            ),
        )


# ============================================================
# OFFICE → PDF
# ============================================================


def office_to_pdf(
    source: Path,
    outdir: Path,
    profile: Path,
):

    binary = binary_path(
        "soffice",
        "libreoffice",
    )

    outdir.mkdir(
        parents=True,
        exist_ok=True,
    )

    profile.mkdir(
        parents=True,
        exist_ok=True,
    )

    # LibreOffice requires a unique user profile
    # for concurrent headless conversions.
    profile_uri = (
        "file://"
        + profile.resolve().as_posix()
    )

    args = [
        binary,
        "--headless",
        "--invisible",
        "--nodefault",
        "--nofirststartwizard",
        "--nologo",
        "--convert-to",
        "pdf",
        "--outdir",
        str(outdir),
        f"-env:UserInstallation={profile_uri}",
        str(source),
    ]

    run_checked(
        args,
        timeout=240,
    )

    output = (
        outdir
        / f"{source.stem}.pdf"
    )

    return verify_output(
        output,
        "Office → PDF",
    )


# ============================================================
# HTML → PDF
# ============================================================


def html_to_pdf(
    source: Path,
    outdir: Path,
    profile: Path,
):

    # LibreOffice is used as the deterministic
    # server-side renderer.
    #
    # For HTML requiring full Chromium CSS/JS,
    # Playwright can be added separately.

    return office_to_pdf(
        source,
        outdir,
        profile,
    )


# ============================================================
# QPDF
# ============================================================


def qpdf_transform(
    source: Path,
    output: Path,
    password: str | None,
    mode: str,
):

    qpdf = binary_path(
        "qpdf"
    )

    output.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    if mode == "unlock":

        args = [
            qpdf,
            "--password="
            + (
                password
                or ""
            ),
            "--decrypt",
            str(source),
            str(output),
        ]

    elif mode == "protect":

        if not password:

            raise HTTPException(
                status_code=400,
                detail=(
                    "A password is required "
                    "to protect the PDF."
                ),
            )

        args = [
            qpdf,
            "--encrypt",
            password,
            password,
            "256",
            "--",
            str(source),
            str(output),
        ]

    else:

        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported qpdf operation."
            ),
        )

    run_checked(
        args,
        timeout=240,
    )

    return verify_output(
        output,
        "PDF security operation",
    )


# ============================================================
# PDF/A
# ============================================================


def find_ghostscript_icc() -> str | None:

    candidates = [
        os.getenv(
            "GS_ICC_PROFILE"
        ),

        "/usr/share/color/icc/ghostscript/srgb.icc",

        "/usr/share/ghostscript/iccprofiles/srgb.icc",
    ]

    if os.name == "nt":

        candidates.extend(
            [
                r"C:\Program Files\gs\gs10.00.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.01.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.02.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.03.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.04.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.05.0\iccprofiles\srgb.icc",
            ]
        )

    for candidate in candidates:

        if candidate and Path(
            candidate
        ).exists():

            return candidate

    return None


def pdf_to_pdfa(
    source: Path,
    output: Path,
):

    gs = binary_path(
        "gs",
        "gswin64c",
        "gswin32c",
    )

    icc_profile = (
        find_ghostscript_icc()
    )

    if not icc_profile:

        raise HTTPException(
            status_code=503,
            detail=(
                "Ghostscript is installed, but "
                "an sRGB ICC profile could not be "
                "found for PDF/A conversion."
            ),
        )

    output.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    args = [
        gs,
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dSAFER",
        "-sDEVICE=pdfwrite",
        "-sColorConversionStrategy=RGB",
        "-sProcessColorModel=DeviceRGB",
        f"-sOutputICCProfile={icc_profile}",
        "-dPDFACompatibilityPolicy=1",
        "-dAutoRotatePages=/None",
        "-o",
        str(output),
        str(source),
    ]

    run_checked(
        args,
        timeout=300,
    )

    return verify_output(
        output,
        "PDF/A conversion",
    )


# ============================================================
# PDF TEXT
# ============================================================


def extract_pdf_text(
    source: Path,
) -> str:

    import fitz

    pdf = _fitz_open(
        source
    )

    try:

        pages = []

        for page in pdf:

            text = (
                page.get_text(
                    "text"
                )
                or ""
            )

            pages.append(
                text.strip()
            )

        return "\n\n".join(
            page
            for page in pages
            if page
        )

    finally:

        pdf.close()


# ============================================================
# TRANSLATION
# ============================================================


# ============================================================
# TRANSLATION ENGINE
# ============================================================

def translation_request(
    text: str,
    target_lang: str,
) -> str:
    """
    Translate text using one of these providers:

    1. Google Cloud Translation v2
       GOOGLE_TRANSLATE_API_KEY

    2. Generic JSON translation endpoint
       TRANSLATION_API_URL

    The Google provider is preferred because the PDF
    reconstruction code only needs a normal translated string.
    """

    text = (text or "").strip()
    target_lang = (target_lang or "").strip()

    if not text:
        return ""

    if not target_lang:
        raise HTTPException(
            status_code=422,
            detail="Target language is required.",
        )

    # --------------------------------------------------------
    # Provider 1: Google Cloud Translation
    # --------------------------------------------------------

    google_key = os.getenv(
        "GOOGLE_TRANSLATE_API_KEY",
        "",
    ).strip()

    if google_key:
        endpoint = (
            "https://translation.googleapis.com/"
            "language/translate/v2"
        )

        payload = json.dumps(
            {
                "q": text,
                "target": target_lang,
                "format": "text",
            }
        ).encode("utf-8")

        request = urllib.request.Request(
            f"{endpoint}?key={quote(google_key)}",
            data=payload,
            headers={
                "Content-Type": "application/json",
                "Accept": "application/json",
            },
            method="POST",
        )

        try:
            with urllib.request.urlopen(
                request,
                timeout=60,
            ) as response:

                body = json.loads(
                    response.read().decode("utf-8")
                )

            translated = (
                body
                .get("data", {})
                .get("translations", [{}])[0]
                .get("translatedText", "")
            )

            if translated:
                return str(translated)

        except Exception as exc:
            raise HTTPException(
                status_code=502,
                detail=(
                    "Google Translation failed: "
                    f"{exc}"
                ),
            )

    # --------------------------------------------------------
    # Provider 2: Generic translation API
    # --------------------------------------------------------

    endpoint = os.getenv(
        "TRANSLATION_API_URL",
        "",
    ).strip()

    if not endpoint:
        raise HTTPException(
            status_code=503,
            detail=(
                "No translation provider is configured. "
                "Set GOOGLE_TRANSLATE_API_KEY or "
                "TRANSLATION_API_URL in server/.env."
            ),
        )

    payload = json.dumps(
        {
            "q": text,
            "target": target_lang,
            "source": "auto",
            "format": "text",
        }
    ).encode("utf-8")

    request = urllib.request.Request(
        endpoint,
        data=payload,
        headers={
            "Content-Type": "application/json",
            "Accept": "application/json",
        },
        method="POST",
    )

    try:
        with urllib.request.urlopen(
            request,
            timeout=60,
        ) as response:

            body = json.loads(
                response.read().decode("utf-8")
            )

    except Exception as exc:
        raise HTTPException(
            status_code=502,
            detail=(
                f"Translation service failed: {exc}"
            ),
        )

    translated = (
        body.get("translatedText")
        or body.get("translation")
        or body.get("text")
        or ""
    )

    if not translated:
        raise HTTPException(
            status_code=502,
            detail=(
                "Translation provider returned "
                "no translated text."
            ),
        )

    return str(translated)

# ============================================================
# SEND EMAIL
# ============================================================


@app.post(
    "/send-email"
)
async def send_email(
    file: UploadFile = File(...),
    to: str = Form(...),
    subject: str = Form(...),
    tool: str = Form(...),
):

    api_key = os.getenv(
        "RESEND_API_KEY"
    )

    email_from = os.getenv(
        "EMAIL_FROM",
        "QuadraConverter <onboarding@resend.dev>",
    )

    if not api_key:

        raise HTTPException(
            status_code=503,
            detail=(
                "Email service is not configured."
            ),
        )

    if (
        not to
        or "@"
        not in to
    ):

        raise HTTPException(
            status_code=400,
            detail=(
                "Please provide a valid email address."
            ),
        )

    content = await file.read()

    if not content:

        raise HTTPException(
            status_code=400,
            detail=(
                "The converted file is empty."
            ),
        )

    if (
        len(content)
        > MAX_FILE_BYTES
    ):

        raise HTTPException(
            status_code=413,
            detail=(
                "The file is too large to email."
            ),
        )

    try:

        import resend

        resend.api_key = (
            api_key
        )

        attachment = {
            "filename": (
                file.filename
                or "converted-file"
            ),
            "content": list(
                content
            ),
        }

        response = (
            resend.Emails.send(
                {
                    "from": email_from,
                    "to": [to],
                    "subject": subject,
                    "html": f"""
                        <div style="font-family:Arial,sans-serif">
                            <h2>QuadraConverter</h2>

                            <p>
                                Your file has been converted
                                successfully using
                                <strong>{tool}</strong>.
                            </p>

                            <p>
                                The converted file is attached
                                to this email.
                            </p>

                            <p>
                                — QuadraConverter
                            </p>
                        </div>
                    """,
                    "attachments": [
                        attachment
                    ],
                }
            )
        )

        return JSONResponse(
            {
                "success": True,
                "message": (
                    "Email sent successfully."
                ),
                "id": (
                    response.get("id")
                    if isinstance(
                        response,
                        dict,
                    )
                    else None
                ),
            }
        )

    except HTTPException:
        raise

    except Exception as exc:

        raise HTTPException(
            status_code=502,
            detail=(
                f"Email provider failed: "
                f"{exc}"
            ),
        )


# ============================================================
# HEALTH CHECK
# ============================================================


@app.get(
    "/health"
)
def health():

    engines = []

    checks = [
        (
            "LibreOffice",
            (
                "soffice",
                "libreoffice",
            ),
        ),
        (
            "qpdf",
            (
                "qpdf",
            ),
        ),
        (
            "Ghostscript",
            (
                "gs",
                "gswin64c",
                "gswin32c",
            ),
        ),
        (
            "Tesseract",
            (
                "tesseract",
            ),
        ),
    ]

    for label, names in checks:

        try:

            path = binary_path(
                *names
            )

            engines.append(
                {
                    "name": label,
                    "available": True,
                    "binary": path,
                }
            )

        except Exception:

            engines.append(
                {
                    "name": label,
                    "available": False,
                    "binary": None,
                }
            )

    return {
        "ok": True,
        "service": APP_NAME,
        "version": APP_VERSION,
        "max_file_mb": (
            MAX_FILE_BYTES
            // (1024 * 1024)
        ),
        "ocr_dpi": OCR_DPI,
        "engines": engines,
    }



def _ensure_signing_identity() -> tuple[str, str]:
    """Return a signing key/certificate pair.

    Production deployments should point these settings at a CA-issued
    certificate and private key. For a ready-to-run install, generate a
    persistent self-signed certificate so the PDF is still cryptographically
    signed and verifiable. Adobe will show a trust warning until the CA chain
    is trusted by the reader.
    """
    if SIGNING_KEY_FILE and SIGNING_CERT_FILE:
        return SIGNING_KEY_FILE, SIGNING_CERT_FILE

    from cryptography import x509
    from cryptography.hazmat.primitives import hashes, serialization
    from cryptography.hazmat.primitives.asymmetric import rsa
    from cryptography.x509.oid import NameOID
    from datetime import datetime, timedelta, timezone

    key_dir = Path(tempfile.gettempdir()) / "quadraconverter-signing"
    key_dir.mkdir(parents=True, exist_ok=True)
    key_file = key_dir / "signer-key.pem"
    cert_file = key_dir / "signer-cert.pem"
    if key_file.exists() and cert_file.exists():
        return str(key_file), str(cert_file)

    key = rsa.generate_private_key(public_exponent=65537, key_size=2048)
    subject = issuer = x509.Name([
        x509.NameAttribute(NameOID.ORGANIZATION_NAME, "QuadraConverter"),
        x509.NameAttribute(NameOID.COMMON_NAME, SIGNING_SUBJECT),
    ])
    now = datetime.now(timezone.utc)
    cert = (
        x509.CertificateBuilder()
        .subject_name(subject)
        .issuer_name(issuer)
        .public_key(key.public_key())
        .serial_number(x509.random_serial_number())
        .not_valid_before(now - timedelta(minutes=5))
        .not_valid_after(now + timedelta(days=825))
        .add_extension(x509.BasicConstraints(ca=False, path_length=None), critical=True)
        .add_extension(x509.KeyUsage(digital_signature=True, content_commitment=True, key_encipherment=False, data_encipherment=False, key_agreement=False, key_cert_sign=False, crl_sign=False, encipher_only=False, decipher_only=False), critical=True)
        .sign(key, hashes.SHA256())
    )
    key_file.write_bytes(key.private_bytes(serialization.Encoding.PEM, serialization.PrivateFormat.PKCS8, serialization.NoEncryption()))
    try:
        os.chmod(key_file, 0o600)
    except OSError:
        pass
    cert_file.write_bytes(cert.public_bytes(serialization.Encoding.PEM))
    return str(key_file), str(cert_file)


def _pki_sign_pdf(source: bytes, signer_name: str, page: int, x: float, y: float, width: float) -> bytes:
    from io import BytesIO
    from pyhanko import stamp
    from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
    from pyhanko.sign import fields, signers, timestamps

    key_file, cert_file = _ensure_signing_identity()
    chain = (SIGNING_CHAIN_FILE,) if SIGNING_CHAIN_FILE else None
    signer = signers.SimpleSigner.load(
        key_file,
        cert_file,
        ca_chain_files=chain,
        key_passphrase=SIGNING_KEY_PASSPHRASE.encode() if SIGNING_KEY_PASSPHRASE else None,
    )

    # Translate the UI's top-left percentage coordinates to PDF's lower-left
    # coordinate system for the visible cryptographic signature field.
    import fitz
    doc = fitz.open(stream=source, filetype="pdf")
    page_index = max(0, min(len(doc) - 1, int(page) - 1))
    rect = doc[page_index].rect
    page_width = float(rect.width)
    page_height = float(rect.height)
    box_width = max(120.0, page_width * max(0.08, min(0.6, float(width) / 100.0)))
    box_height = max(55.0, min(95.0, box_width * 0.34))
    ll_x = max(0.0, min(page_width - box_width, page_width * float(x) / 100.0))
    top_y = page_height * float(y) / 100.0
    ll_y = max(0.0, min(page_height - box_height, page_height - top_y - box_height))
    doc.close()

    writer = IncrementalPdfFileWriter(BytesIO(source))
    field_name = f"QuadraSignature_{os.urandom(6).hex()}"
    new_field = fields.SigFieldSpec(field_name, on_page=page_index, box=(int(ll_x), int(ll_y), int(ll_x + box_width), int(ll_y + box_height)))
    meta = signers.PdfSignatureMetadata(
        field_name=field_name,
        md_algorithm="sha256",
        name=signer_name or SIGNING_SUBJECT,
        reason="Electronic signature",
        subfilter=fields.SigSeedSubFilter.PADES,
    )
    timestamper = timestamps.HTTPTimeStamper(SIGNING_TSA_URL) if SIGNING_TSA_URL else None
    pdf_signer = signers.PdfSigner(
        meta,
        signer=signer,
        timestamper=timestamper,
        stamp_style=stamp.TextStampStyle(stamp_text="Digitally signed by %(signer)s\nTimestamp: %(ts)s"),
        new_field_spec=new_field,
    )
    output = BytesIO()
    pdf_signer.sign_pdf(writer, output=output)
    return output.getvalue()


@app.post("/sign-pdf")
async def sign_pdf_endpoint(
    file: UploadFile = File(...),
    signer_name: str = Form("QuadraConverter User"),
    page: int = Form(1),
    x: float = Form(68),
    y: float = Form(74),
    width: float = Form(24),
):
    """Embed a real PDF cryptographic signature using X.509/PKI material."""
    work = Path(tempfile.mkdtemp(prefix="quadra-sign-"))
    try:
        data = await file.read()
        if not data or len(data) > MAX_FILE_BYTES:
            raise HTTPException(status_code=400, detail="Invalid or oversized PDF.")
        if not data.startswith(b"%PDF-"):
            raise HTTPException(status_code=400, detail="The signing endpoint accepts PDF files only.")
        output = _pki_sign_pdf(data, signer_name, page, x, y, width)
        filename = safe_filename(file.filename, "document.pdf").rsplit('.', 1)[0] + "-pki-signed.pdf"
        return StreamingResponse(BytesIO(output), media_type="application/pdf", headers={"Content-Disposition": f'attachment; filename="{filename}"', "X-Signing-Mode": "PAdES-PKI"})
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"PKI signing failed: {exc}") from exc
    finally:
        cleanup(work)


# ============================================================
# SINGLE CONVERSION ENDPOINT
#
# DO NOT CREATE ANOTHER @app.post("/convert") BELOW THIS.
# ============================================================
# ============================================================
# QUADRA AI PDF ENGINE
# ============================================================

AI_BASE_URL = os.getenv(
    "AI_BASE_URL",
    "https://openrouter.ai/api/v1",
).rstrip("/")

AI_API_KEY = os.getenv(
    "OPENROUTER_API_KEY",
    "",
).strip()

AI_CHAT_MODEL = os.getenv(
    "AI_CHAT_MODEL",
    "openai/gpt-4o-mini",
).strip()

AI_TRANSLATION_MODEL = os.getenv(
    "AI_TRANSLATION_MODEL",
    "openai/gpt-4o-mini",
).strip()

AI_EMBEDDING_MODEL = os.getenv(
    "AI_EMBEDDING_MODEL",
    "openai/text-embedding-3-small",
).strip()

SUPABASE_URL = os.getenv(
    "SUPABASE_URL",
    "",
).rstrip("/")

SUPABASE_SERVICE_ROLE_KEY = os.getenv(
    "SUPABASE_SERVICE_ROLE_KEY",
    "",
).strip()

PDF_CHUNK_SIZE = int(
    os.getenv(
        "PDF_CHUNK_SIZE",
        "1200",
    )
)

PDF_CHUNK_OVERLAP = int(
    os.getenv(
        "PDF_CHUNK_OVERLAP",
        "200",
    )
)

PDF_RAG_TOP_K = int(
    os.getenv(
        "PDF_RAG_TOP_K",
        "6",
    )
)


def require_ai_config():

    if not AI_API_KEY:
        raise HTTPException(
            status_code=503,
            detail=(
                "OPENROUTER_API_KEY is not configured "
                "on the conversion server."
            ),
        )


def require_supabase_config():

    if not SUPABASE_URL:
        raise HTTPException(
            status_code=503,
            detail=(
                "SUPABASE_URL is not configured "
                "on the conversion server."
            ),
        )

    if not SUPABASE_SERVICE_ROLE_KEY:
        raise HTTPException(
            status_code=503,
            detail=(
                "SUPABASE_SERVICE_ROLE_KEY is not configured "
                "on the conversion server."
            ),
        )


def supabase_headers():

    require_supabase_config()

    return {
        "apikey": SUPABASE_SERVICE_ROLE_KEY,
        "Authorization": (
            f"Bearer {SUPABASE_SERVICE_ROLE_KEY}"
        ),
        "Content-Type": "application/json",
        "Prefer": "return=representation",
    }


def sha256_file(
    path: Path,
) -> str:

    digest = hashlib.sha256()

    with path.open("rb") as handle:

        while True:

            chunk = handle.read(
                1024 * 1024
            )

            if not chunk:
                break

            digest.update(chunk)

    return digest.hexdigest()


def chunk_text(
    text: str,
    chunk_size: int = PDF_CHUNK_SIZE,
    overlap: int = PDF_CHUNK_OVERLAP,
) -> list[str]:

    cleaned = re.sub(
        r"[ \t]+",
        " ",
        text,
    )

    cleaned = re.sub(
        r"\n{3,}",
        "\n\n",
        cleaned,
    ).strip()

    if not cleaned:
        return []

    words = cleaned.split()

    if not words:
        return []

    chunks: list[str] = []

    current: list[str] = []
    current_length = 0

    for word in words:

        word_length = len(word) + 1

        if (
            current
            and current_length + word_length > chunk_size
        ):

            chunks.append(
                " ".join(current).strip()
            )

            overlap_words: list[str] = []
            overlap_length = 0

            for previous in reversed(current):

                if (
                    overlap_length
                    + len(previous)
                    + 1
                    > overlap
                ):
                    break

                overlap_words.insert(
                    0,
                    previous,
                )

                overlap_length += (
                    len(previous) + 1
                )

            current = overlap_words
            current_length = overlap_length

        current.append(word)
        current_length += word_length

    if current:
        chunks.append(
            " ".join(current).strip()
        )

    return [
        value
        for value in chunks
        if len(value.strip()) >= 20
    ]


def extract_pdf_blocks_for_ai(
    source: Path,
    ocr_language: str = "eng",
) -> list[dict[str, Any]]:

    import fitz

    pdf = _fitz_open(
        source
    )

    blocks: list[dict[str, Any]] = []

    try:

        for page_index, page in enumerate(
            pdf,
            start=1,
        ):

            native_blocks = page.get_text(
                "blocks"
            )

            native_text = (
                page.get_text(
                    "text"
                )
                or ""
            ).strip()

            useful_native = [
                block
                for block in native_blocks
                if len(
                    str(
                        block[4]
                        if len(block) > 4
                        else ""
                    ).strip()
                ) >= 2
            ]

            if len(native_text) >= OCR_MIN_TEXT_CHARS:

                for block in useful_native:

                    text = str(
                        block[4]
                    ).strip()

                    if not text:
                        continue

                    blocks.append(
                        {
                            "page": page_index,
                            "text": text,
                            "bbox": [
                                float(block[0]),
                                float(block[1]),
                                float(block[2]),
                                float(block[3]),
                            ],
                            "source": "native",
                        }
                    )

                continue

            if not tesseract_available():

                continue

            try:

                pixmap, _ = render_page_image(
                    page,
                    OCR_DPI,
                )

                image = Image.open(
                    BytesIO(
                        pixmap.tobytes(
                            "png"
                        )
                    )
                )

                data = pytesseract.image_to_data(
                    image,
                    lang=normalize_ocr_language(
                        ocr_language
                    ),
                    output_type=(
                        pytesseract.Output.DICT
                    ),
                )

                scale = 72.0 / float(
                    OCR_DPI
                )

                grouped: dict[
                    tuple[int, int, int],
                    list[tuple[str, int, int, int, int]]
                ] = {}

                count = len(
                    data.get(
                        "text",
                        [],
                    )
                )

                for i in range(count):

                    value = str(
                        data["text"][i]
                        or ""
                    ).strip()

                    if not value:
                        continue

                    confidence = str(
                        data.get(
                            "conf",
                            ["-1"] * count,
                        )[i]
                    )

                    try:
                        if float(confidence) < 20:
                            continue
                    except Exception:
                        pass

                    key = (
                        int(
                            data["block_num"][i]
                        ),
                        int(
                            data["par_num"][i]
                        ),
                        int(
                            data["line_num"][i]
                        ),
                    )

                    grouped.setdefault(
                        key,
                        [],
                    ).append(
                        (
                            value,
                            int(
                                data["left"][i]
                            ),
                            int(
                                data["top"][i]
                            ),
                            int(
                                data["width"][i]
                            ),
                            int(
                                data["height"][i]
                            ),
                        )
                    )

                for words in grouped.values():

                    if not words:
                        continue

                    text = " ".join(
                        item[0]
                        for item in words
                    ).strip()

                    if not text:
                        continue

                    left = min(
                        item[1]
                        for item in words
                    )

                    top = min(
                        item[2]
                        for item in words
                    )

                    right = max(
                        item[1] + item[3]
                        for item in words
                    )

                    bottom = max(
                        item[2] + item[4]
                        for item in words
                    )

                    blocks.append(
                        {
                            "page": page_index,
                            "text": text,
                            "bbox": [
                                left * scale,
                                top * scale,
                                right * scale,
                                bottom * scale,
                            ],
                            "source": "ocr",
                        }
                    )

            except Exception as exc:

                print(
                    "[QuadraAI] OCR failed "
                    f"on page {page_index}: {exc}"
                )

    finally:

        pdf.close()

    return blocks


def openrouter_request(
    endpoint: str,
    payload: dict[str, Any],
) -> dict[str, Any]:

    require_ai_config()

    headers = {
        "Authorization": (
            f"Bearer {AI_API_KEY}"
        ),
        "Content-Type": "application/json",
        "Accept": "application/json",
        "HTTP-Referer": os.getenv(
            "AI_HTTP_REFERER",
            "https://quadraconverter.in",
        ),
        "X-Title": "QuadraConverter",
    }

    try:

        response = httpx.post(
            f"{AI_BASE_URL}/{endpoint.lstrip('/')}",
            headers=headers,
            json=payload,
            timeout=180,
        )

    except httpx.TimeoutException:

        raise HTTPException(
            status_code=504,
            detail=(
                "AI service timed out. "
                "Please try again."
            ),
        )

    except httpx.HTTPError as exc:

        raise HTTPException(
            status_code=502,
            detail=(
                f"AI service connection failed: {exc}"
            ),
        )

    if response.status_code >= 400:

        try:
            body = response.json()
            detail = (
                body.get("error", {}).get("message")
                if isinstance(
                    body.get("error"),
                    dict,
                )
                else body.get("error")
            )
        except Exception:
            detail = response.text

        raise HTTPException(
            status_code=502,
            detail=(
                str(detail)
                or "AI service returned an error."
            ),
        )

    try:
        return response.json()
    except Exception:

        raise HTTPException(
            status_code=502,
            detail="AI service returned invalid JSON.",
        )


def create_embedding(
    text: str,
) -> list[float]:

    payload = {
        "model": AI_EMBEDDING_MODEL,
        "input": text,
    }

    body = openrouter_request(
        "/embeddings",
        payload,
    )

    data = body.get(
        "data"
    )

    if not isinstance(
        data,
        list,
    ) or not data:

        raise HTTPException(
            status_code=502,
            detail="Embedding service returned no vector.",
        )

    embedding = data[0].get(
        "embedding"
    )

    if not isinstance(
        embedding,
        list,
    ):

        raise HTTPException(
            status_code=502,
            detail="Embedding service returned an invalid vector.",
        )

    return [
        float(value)
        for value in embedding
    ]


def ai_chat_completion(
    messages: list[dict[str, str]],
    model: str = AI_CHAT_MODEL,
    temperature: float = 0.1,
) -> str:

    payload = {
        "model": model,
        "messages": messages,
        "temperature": temperature,
    }

    body = openrouter_request(
        "/chat/completions",
        payload,
    )

    choices = body.get(
        "choices"
    )

    if not choices:
        raise HTTPException(
            status_code=502,
            detail="AI model returned no answer.",
        )

    message = (
        choices[0]
        .get("message", {})
    )

    content = message.get(
        "content"
    )

    if not isinstance(
        content,
        str,
    ):

        raise HTTPException(
            status_code=502,
            detail="AI model returned an invalid answer.",
        )

    return content.strip()


def supabase_get_document_by_hash(
    file_hash: str,
) -> dict[str, Any] | None:

    require_supabase_config()

    response = httpx.get(
        f"{SUPABASE_URL}/rest/v1/pdf_documents",
        headers=supabase_headers(),
        params={
            "file_hash": f"eq.{file_hash}",
            "select": "id,file_hash,file_name,page_count",
            "limit": "1",
        },
        timeout=60,
    )

    if response.status_code >= 400:

        raise HTTPException(
            status_code=502,
            detail=(
                "Could not read PDF document index "
                f"from Supabase: {response.text}"
            ),
        )

    rows = response.json()

    return rows[0] if rows else None


def supabase_create_document(
    file_hash: str,
    file_name: str,
    page_count: int,
) -> dict[str, Any]:

    require_supabase_config()

    response = httpx.post(
        f"{SUPABASE_URL}/rest/v1/pdf_documents",
        headers=supabase_headers(),
        json={
            "file_hash": file_hash,
            "file_name": file_name,
            "page_count": page_count,
        },
        timeout=60,
    )

    if response.status_code >= 400:

        if response.status_code == 409:

            existing = (
                supabase_get_document_by_hash(
                    file_hash
                )
            )

            if existing:
                return existing

        raise HTTPException(
            status_code=502,
            detail=(
                "Could not create PDF document "
                f"in Supabase: {response.text}"
            ),
        )

    rows = response.json()

    if not rows:
        raise HTTPException(
            status_code=502,
            detail="Supabase did not return the PDF document.",
        )

    return rows[0]


def supabase_insert_chunks(
    rows: list[dict[str, Any]],
):

    if not rows:
        return

    response = httpx.post(
        f"{SUPABASE_URL}/rest/v1/pdf_chunks",
        headers=supabase_headers(),
        json=rows,
        timeout=180,
    )

    if response.status_code >= 400:

        raise HTTPException(
            status_code=502,
            detail=(
                "Could not store PDF embeddings "
                f"in Supabase: {response.text}"
            ),
        )


def ensure_pdf_indexed(
    source: Path,
    original_filename: str,
    ocr_language: str = "eng",
) -> tuple[str, int]:

    file_hash = sha256_file(
        source
    )

    existing = (
        supabase_get_document_by_hash(
            file_hash
        )
    )

    if existing:

        return (
            str(existing["id"]),
            int(
                existing.get(
                    "page_count",
                    0,
                )
            ),
        )

    blocks = extract_pdf_blocks_for_ai(
        source,
        ocr_language,
    )

    if not blocks:

        raise HTTPException(
            status_code=422,
            detail=(
                "No readable text was found in this PDF. "
                "The document may contain unsupported scans."
            ),
        )

    import fitz

    pdf = _fitz_open(
        source
    )

    try:
        page_count = pdf.page_count
    finally:
        pdf.close()

    document = (
        supabase_create_document(
            file_hash,
            safe_filename(
                original_filename,
                "document.pdf",
            ),
            page_count,
        )
    )

    document_id = str(
        document["id"]
    )

    chunks: list[dict[str, Any]] = []

    chunk_index = 0

    for block in blocks:

        pieces = chunk_text(
            block["text"]
        )

        for piece in pieces:

            embedding = create_embedding(
                piece
            )

            chunks.append(
                {
                    "document_id": document_id,
                    "chunk_index": chunk_index,
                    "page_number": int(
                        block["page"]
                    ),
                    "content": piece,
                    "embedding": embedding,
                }
            )

            chunk_index += 1

            if len(chunks) >= 20:

                supabase_insert_chunks(
                    chunks
                )

                chunks = []

    if chunks:
        supabase_insert_chunks(
            chunks
        )

    return (
        document_id,
        page_count,
    )


def supabase_match_chunks(
    document_id: str,
    query_embedding: list[float],
    top_k: int = PDF_RAG_TOP_K,
) -> list[dict[str, Any]]:

    require_supabase_config()

    response = httpx.post(
        f"{SUPABASE_URL}/rest/v1/rpc/match_pdf_chunks",
        headers=supabase_headers(),
        json={
            "query_embedding": query_embedding,
            "match_document_id": document_id,
            "match_count": top_k,
            "min_similarity": 0.12,
        },
        timeout=120,
    )

    if response.status_code >= 400:

        raise HTTPException(
            status_code=502,
            detail=(
                "Vector search failed: "
                f"{response.text}"
            ),
        )

    result = response.json()

    if not isinstance(
        result,
        list,
    ):
        return []

    return result


def chat_with_pdf_document(
    source: Path,
    original_filename: str,
    question: str,
    ocr_language: str = "eng",
) -> str:

    question = question.strip()

    if not question:

        raise HTTPException(
            status_code=400,
            detail="Please enter a question about the PDF.",
        )

    document_id, _ = ensure_pdf_indexed(
        source,
        original_filename,
        ocr_language,
    )

    question_embedding = create_embedding(
        question
    )

    matches = supabase_match_chunks(
        document_id,
        question_embedding,
        PDF_RAG_TOP_K,
    )

    if not matches:

        return (
            "I could not find relevant information "
            "in the uploaded PDF."
        )

    context_parts: list[str] = []

    for index, match in enumerate(
        matches,
        start=1,
    ):

        page = match.get(
            "page_number"
        )

        content = str(
            match.get(
                "content",
                "",
            )
        ).strip()

        if not content:
            continue

        context_parts.append(
            (
                f"[Source {index} | "
                f"Page {page or '?'}]\n"
                f"{content}"
            )
        )

    context = "\n\n".join(
        context_parts
    )

    system_prompt = """
You are QuadraConverter PDF Assistant.

Answer questions ONLY from the supplied PDF context.

Rules:
1. Do not invent information.
2. If the answer is not contained in the context, say that it was not found in the PDF.
3. Be clear and concise.
4. Preserve important numbers, dates, names and technical terminology.
5. Mention the relevant page number when possible.
6. Do not use outside knowledge to fill missing information.
"""

    user_prompt = (
        "PDF CONTEXT:\n\n"
        f"{context}\n\n"
        "QUESTION:\n"
        f"{question}"
    )

    return ai_chat_completion(
        [
            {
                "role": "system",
                "content": system_prompt.strip(),
            },
            {
                "role": "user",
                "content": user_prompt,
            },
        ],
        model=AI_CHAT_MODEL,
        temperature=0.05,
    )


# ============================================================
# PDF TRANSLATION WITH LAYOUT PRESERVATION
# ============================================================

INDIC_LANGUAGE_NAMES = {
    "english": "English",
    "en": "English",

    "hindi": "Hindi",
    "hi": "Hindi",

    "tamil": "Tamil",
    "ta": "Tamil",

    "telugu": "Telugu",
    "te": "Telugu",

    "malayalam": "Malayalam",
    "ml": "Malayalam",

    "kannada": "Kannada",
    "kn": "Kannada",

    "marathi": "Marathi",
    "mr": "Marathi",

    "bengali": "Bengali",
    "bn": "Bengali",

    "gujarati": "Gujarati",
    "gu": "Gujarati",

    "punjabi": "Punjabi",
    "pa": "Punjabi",

    "odia": "Odia",
    "or": "Odia",

    "assamese": "Assamese",
    "as": "Assamese",

    "urdu": "Urdu",
    "ur": "Urdu",

    "nepali": "Nepali",
    "ne": "Nepali",

    "sanskrit": "Sanskrit",
    "sa": "Sanskrit",

    "french": "French",
    "fr": "French",

    "german": "German",
    "de": "German",

    "spanish": "Spanish",
    "es": "Spanish",

    "italian": "Italian",
    "it": "Italian",

    "portuguese": "Portuguese",
    "pt": "Portuguese",
}


def language_display_name(
    language: str,
) -> str:

    value = (
        language
        or "English"
    ).strip().lower()

    return INDIC_LANGUAGE_NAMES.get(
        value,
        language.strip() or "English",
    )


def translate_text_ai(
    text: str,
    target_language: str,
) -> str:

    target = language_display_name(
        target_language
    )

    source_language = (
        "auto-detect"
    )

    prompt = f"""
Translate the following document text into {target}.

Source language: {source_language}

Requirements:
- Preserve the exact meaning.
- Preserve numbers, dates, names, units and technical terms.
- Do not summarize.
- Do not add explanations.
- Keep paragraph structure.
- Do not translate proper brand/product names unless appropriate.
- For Indian languages, use natural native grammar and terminology.
- Return only the translation.

TEXT:
{text}
"""

    return ai_chat_completion(
        [
            {
                "role": "system",
                "content": (
                    "You are a professional document "
                    "translation engine."
                ),
            },
            {
                "role": "user",
                "content": prompt.strip(),
            },
        ],
        model=AI_TRANSLATION_MODEL,
        temperature=0.05,
    )


def find_pdf_font(
    target_language: str,
) -> str | None:

    language = (
        target_language or ""
    ).strip().lower()

    font_names = {
        "tamil": "NotoSansTamil-Regular.ttf",
        "ta": "NotoSansTamil-Regular.ttf",

        "hindi": "NotoSansDevanagari-Regular.ttf",
        "hi": "NotoSansDevanagari-Regular.ttf",

        "marathi": "NotoSansDevanagari-Regular.ttf",
        "mr": "NotoSansDevanagari-Regular.ttf",

        "nepali": "NotoSansDevanagari-Regular.ttf",
        "ne": "NotoSansDevanagari-Regular.ttf",

        "sanskrit": "NotoSansDevanagari-Regular.ttf",
        "sa": "NotoSansDevanagari-Regular.ttf",

        "telugu": "NotoSansTelugu-Regular.ttf",
        "te": "NotoSansTelugu-Regular.ttf",

        "kannada": "NotoSansKannada-Regular.ttf",
        "kn": "NotoSansKannada-Regular.ttf",

        "malayalam": "NotoSansMalayalam-Regular.ttf",
        "ml": "NotoSansMalayalam-Regular.ttf",

        "bengali": "NotoSansBengali-Regular.ttf",
        "bn": "NotoSansBengali-Regular.ttf",

        "assamese": "NotoSansBengali-Regular.ttf",
        "as": "NotoSansBengali-Regular.ttf",

        "gujarati": "NotoSansGujarati-Regular.ttf",
        "gu": "NotoSansGujarati-Regular.ttf",

        "punjabi": "NotoSansGurmukhi-Regular.ttf",
        "pa": "NotoSansGurmukhi-Regular.ttf",

        "odia": "NotoSansOriya-Regular.ttf",
        "or": "NotoSansOriya-Regular.ttf",

        "urdu": "NotoNaskhArabic-Regular.ttf",
        "ur": "NotoNaskhArabic-Regular.ttf",
    }

    wanted = font_names.get(
        language,
        "NotoSans-Regular.ttf",
    )

    font_directories = [
        "/usr/share/fonts",
        "/usr/local/share/fonts",
    ]

    for directory in font_directories:

        root = Path(directory)

        if not root.exists():
            continue

        matches = list(
            root.rglob(wanted)
        )

        if matches:
            return str(
                matches[0]
            )

    fallback_names = [
        "NotoSans-Regular.ttf",
        "DejaVuSans.ttf",
    ]

    for name in fallback_names:

        for directory in font_directories:

            root = Path(directory)

            if not root.exists():
                continue

            matches = list(
                root.rglob(name)
            )

            if matches:
                return str(
                    matches[0]
                )

    return None

def translate_pdf_with_layout(
    source: Path,
    output: Path,
    target_language: str,
    ocr_language: str = "eng",
):

    import fitz

    pdf = _fitz_open(
        source
    )

    try:

        font_file = find_pdf_font(
            target_language
        )

        for page in pdf:

            blocks = page.get_text(
                "blocks"
            )

            native_text = (
                page.get_text(
                    "text"
                )
                or ""
            ).strip()

            if len(native_text) < OCR_MIN_TEXT_CHARS:

                if tesseract_available():

                    try:

                        pixmap, _ = render_page_image(
                            page,
                            OCR_DPI,
                        )

                        image = Image.open(
                            BytesIO(
                                pixmap.tobytes(
                                    "png"
                                )
                            )
                        )

                        data = pytesseract.image_to_data(
                            image,
                            lang=normalize_ocr_language(
                                ocr_language
                            ),
                            output_type=(
                                pytesseract.Output.DICT
                            ),
                        )

                        scale = 72.0 / float(
                            OCR_DPI
                        )

                        for i in range(
                            len(
                                data.get(
                                    "text",
                                    [],
                                )
                            )
                        ):

                            value = str(
                                data["text"][i]
                                or ""
                            ).strip()

                            if not value:
                                continue

                            left = (
                                float(
                                    data["left"][i]
                                )
                                * scale
                            )

                            top = (
                                float(
                                    data["top"][i]
                                )
                                * scale
                            )

                            width = (
                                float(
                                    data["width"][i]
                                )
                                * scale
                            )

                            height = (
                                float(
                                    data["height"][i]
                                )
                                * scale
                            )

                            rect = fitz.Rect(
                                left,
                                top,
                                left + width,
                                top + height,
                            )

                            translated = (
                                translate_text_ai(
                                    value,
                                    target_language,
                                )
                            )

                            page.add_redact_annot(
                                rect,
                                fill=(1, 1, 1),
                            )

                            page.apply_redactions()

                            fontsize = max(
                                6,
                                min(
                                    18,
                                    height * 0.72,
                                ),
                            )

                            if font_file:

                                page.insert_textbox(
                                    rect,
                                    translated,
                                    fontfile=font_file,
                                    fontsize=fontsize,
                                    color=(0, 0, 0),
                                    align=0,
                                )

                            else:

                                page.insert_textbox(
                                    rect,
                                    translated,
                                    fontsize=fontsize,
                                    color=(0, 0, 0),
                                    align=0,
                                )

                    except Exception as exc:

                        print(
                            "[QuadraAI] OCR translation "
                            f"failed: {exc}"
                        )

                continue

            for block in blocks:

                if len(block) < 5:
                    continue

                original = str(
                    block[4]
                ).strip()

                if not original:
                    continue

                rect = fitz.Rect(
                    block[0],
                    block[1],
                    block[2],
                    block[3],
                )

                if rect.width <= 1 or rect.height <= 1:
                    continue

                translated = (
                    translate_text_ai(
                        original,
                        target_language,
                    )
                )

                if not translated:
                    continue

                page.add_redact_annot(
                    rect,
                    fill=(1, 1, 1),
                )

            page.apply_redactions()

            for block in blocks:

                if len(block) < 5:
                    continue

                original = str(
                    block[4]
                ).strip()

                if not original:
                    continue

                rect = fitz.Rect(
                    block[0],
                    block[1],
                    block[2],
                    block[3],
                )

                translated = (
                    translate_text_ai(
                        original,
                        target_language,
                    )
                )

                if not translated:
                    continue

                estimated_font = max(
                    6,
                    min(
                        18,
                        rect.height * 0.65,
                    ),
                )

                if font_file:

                    page.insert_textbox(
                        rect,
                        translated,
                        fontfile=font_file,
                        fontsize=estimated_font,
                        color=(0, 0, 0),
                        align=0,
                    )

                else:

                    page.insert_textbox(
                        rect,
                        translated,
                        fontsize=estimated_font,
                        color=(0, 0, 0),
                        align=0,
                    )

        pdf.save(
            str(output),
            garbage=4,
            deflate=True,
            clean=True,
        )

    finally:

        pdf.close()


# ============================================================
# PDF AI CHAT ROUTE
# ============================================================


@app.post(
    "/pdf-chat"
)
async def pdf_chat(
    file: UploadFile = File(...),
    question: str = Form(...),
    language: str = Form("eng"),
):

    work = Path(
        tempfile.mkdtemp(
            prefix="quadra-pdf-chat-"
        )
    )

    try:

        source = save_upload(
            file,
            work,
            ALLOWED_PDF,
        )

        answer = chat_with_pdf_document(
            source,
            file.filename or "document.pdf",
            question,
            language,
        )

        output = (
            work
            / "pdf-answer.txt"
        )

        output.write_text(
            answer,
            encoding="utf-8",
        )

        return file_response(
            output,
            "text/plain; charset=utf-8",
            work,
            "Quadra AI RAG + Supabase pgvector",
        )

    except HTTPException:

        cleanup(work)
        raise

    except Exception as exc:

        cleanup(work)

        raise HTTPException(
            status_code=500,
            detail=(
                "PDF chat failed: "
                f"{exc}"
            ),
        )

    finally:

        try:
            await file.close()
        except Exception:
            pass

@app.post(
    "/convert"
)
async def convert(
    file: UploadFile = File(...),
    operation: str = Form(...),
    password: str = Form(""),
    targetLang: str = Form(""),
    language: str = Form("eng"),
):

    work = Path(
        tempfile.mkdtemp(
            prefix="quadra-convert-"
        )
    )

    try:

        if not file.filename:

            raise HTTPException(
                status_code=400,
                detail=(
                    "No file was uploaded."
                ),
            )

        operation = (
            operation
            or ""
        ).strip().lower()

        language = (
            language
            or "eng"
        ).strip()

        # ====================================================
        # OFFICE → PDF
        # ====================================================

        if operation == "office-to-pdf":

            source = save_upload(
                file,
                work,
                ALLOWED_OFFICE,
            )

            outdir = (
                work / "out"
            )

            profile = (
                work / "profile"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            profile.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = office_to_pdf(
                source,
                outdir,
                profile,
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "LibreOffice",
            )

        # ====================================================
        # HTML → PDF
        # ====================================================

        if operation == "html-to-pdf":

            source = save_upload(
                file,
                work,
                ALLOWED_HTML,
            )

            outdir = (
                work / "out"
            )

            profile = (
                work / "profile"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            profile.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = html_to_pdf(
                source,
                outdir,
                profile,
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "LibreOffice",
            )

        # ====================================================
        # PDF UNLOCK
        # ====================================================

        if operation == "pdf-unlock":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = qpdf_transform(
                source,
                (
                    outdir
                    / f"{source.stem}-unlocked.pdf"
                ),
                password,
                "unlock",
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "qpdf",
            )

        # ====================================================
        # PDF PROTECT
        # ====================================================

        if operation == "pdf-protect":

            if not password:

                raise HTTPException(
                    status_code=400,
                    detail=(
                        "A password is required "
                        "to protect the PDF."
                    ),
                )

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = qpdf_transform(
                source,
                (
                    outdir
                    / f"{source.stem}-protected.pdf"
                ),
                password,
                "protect",
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "qpdf",
            )

        # ====================================================
        # PDF → PDF/A
        # ====================================================

        if operation == "pdf-to-pdfa":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = pdf_to_pdfa(
                source,
                (
                    outdir
                    / f"{source.stem}-pdfa.pdf"
                ),
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "Ghostscript",
            )

        # ====================================================
        # PDF → WORD
        #
        # IMPORTANT:
        # This is NOT browser-side conversion.
        #
        # Pipeline:
        #
        # PDF
        #  ↓
        # PyMuPDF native text
        #  ↓
        # coordinate extraction
        #  ↓
        # line reconstruction
        #  ↓
        # scanned-page detection
        #  ↓
        # Tesseract OCR when required
        #  ↓
        # Word paragraph reconstruction
        #  ↓
        # DOCX
        # ====================================================

        if operation == "pdf-to-word":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            ocr_language = (
                normalize_ocr_language(
                    language
                )
            )

            output = pdf_to_docx(
                source,
                (
                    outdir
                    / f"{source.stem}.docx"
                ),
                ocr_language,
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                ),
                work,
                "PyMuPDF + Tesseract + python-docx",
            )

        # ====================================================
        # PDF → EXCEL
        #
        # Pipeline:
        #
        # PDF
        #  ↓
        # native table detection
        #  ↓
        # coordinate-based text extraction
        #  ↓
        # column clustering
        #  ↓
        # OCR fallback for scanned PDF
        #  ↓
        # XLSX
        # ====================================================

        if operation == "pdf-to-xlsx":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            ocr_language = (
                normalize_ocr_language(
                    language
                )
            )

            output = pdf_to_xlsx(
                source,
                (
                    outdir
                    / f"{source.stem}.xlsx"
                ),
                ocr_language,
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "spreadsheetml.sheet"
                ),
                work,
                "pdfplumber + PyMuPDF + Tesseract + openpyxl",
            )

        # ====================================================
        # PDF → POWERPOINT
        # ====================================================

        if operation == "pdf-to-pptx":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = pdf_to_pptx(
                source,
                (
                    outdir
                    / f"{source.stem}.pptx"
                ),
                normalize_ocr_language(language),
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                work,
                "PyMuPDF + python-pptx",
            )

        # ====================================================
        # PDF TRANSLATE
        # ====================================================

                # ====================================================
        # PDF TRANSLATE
        # ====================================================

        if operation == "pdf-translate":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            target_language = (
                targetLang
                or "en"
            ).strip()

            if not target_language:

                raise HTTPException(
                    status_code=400,
                    detail=(
                        "Please select a target language."
                    ),
                )

            output = (
                work
                / f"{source.stem}-translated.pdf"
            )

            translate_pdf_with_layout(
                source,
                output,
                target_language,
                language,
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "Quadra AI + PyMuPDF + OCR",
            )
# ============================================================
# STARTUP DIAGNOSTICS
# ============================================================


@app.on_event(
    "startup"
)
def startup_check():

    print("")
    print(
        "============================================================"
    )
    print(
        " QuadraConverter Conversion API"
    )
    print(
        f" Version: {APP_VERSION}"
    )
    print(
        "============================================================"
    )

    engines = [
        (
            "LibreOffice",
            (
                "soffice",
                "libreoffice",
            ),
        ),
        (
            "qpdf",
            (
                "qpdf",
            ),
        ),
        (
            "Ghostscript",
            (
                "gs",
                "gswin64c",
                "gswin32c",
            ),
        ),
        (
            "Tesseract",
            (
                "tesseract",
            ),
        ),
    ]

    for label, names in engines:

        try:

            path = binary_path(
                *names
            )

            print(
                f"[OK] {label}: {path}"
            )

        except Exception:

            print(
                f"[WARN] {label}: NOT INSTALLED"
            )

    print(
        "============================================================"
    )
    print(
        f"Max file size: "
        f"{MAX_FILE_BYTES // (1024 * 1024)} MB"
    )
    print(
        f"OCR DPI: {OCR_DPI}"
    )
    print(
        f"OCR threshold: "
        f"{OCR_MIN_TEXT_CHARS} characters"
    )
    print(
        "============================================================"
    )
    print("")


# ============================================================
# LOCAL DEVELOPMENT
# ============================================================

if __name__ == "__main__":

    import uvicorn

    uvicorn.run(
        "converter_api:app",
        host="0.0.0.0",
        port=int(
            os.getenv(
                "PORT",
                "8000",
            )
        ),
        reload=False,
    )
